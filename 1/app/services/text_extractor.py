# ============================================
# FILE: app/services/text_extractor.py
# Text Extraction Service - Pipeline Component #2
# Crawl → [EXTRACT] → Chunk → Index → Search
# ============================================
"""
Text Extractor Service

Extracts text content from various file formats:
- PDF (using pypdf or pdfminer)
- DOCX (using python-docx)
- DOC (using antiword or fallback)
- TXT, MD, CSV, JSON (plain text)
- HTML (using BeautifulSoup)

Usage:
    extracted = await TextExtractor.extract("/path/to/file.pdf")
    if extracted:
        print(f"Extracted {extracted.word_count} words")
        print(extracted.text)
"""
import os
import re
import hashlib
from pathlib import Path
from typing import Optional, Dict, Any, List
from dataclasses import dataclass, field
import logging

logger = logging.getLogger(__name__)


@dataclass
class ExtractedContent:
    """Represents extracted content from a file"""
    text: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    page_count: Optional[int] = None
    word_count: int = 0
    char_count: int = 0
    content_hash: str = ""
    extraction_method: str = "unknown"
    
    def __post_init__(self):
        if self.text:
            self.word_count = len(self.text.split())
            self.char_count = len(self.text)
            self.content_hash = hashlib.md5(self.text.encode()).hexdigest()


class TextExtractor:
    """
    Extracts text content from various file types.
    
    Supports: PDF, DOCX, DOC, TXT, MD, RTF, HTML, CSV, JSON, XML
    """
    
    SUPPORTED_EXTENSIONS = {
        '.pdf': 'pdf',
        '.docx': 'docx',
        '.doc': 'doc',
        '.pptx': 'pptx',
        '.ppt': 'ppt',
        '.txt': 'text',
        '.md': 'markdown',
        '.markdown': 'markdown',
        '.rtf': 'rtf',
        '.html': 'html',
        '.htm': 'html',
        '.csv': 'csv',
        '.json': 'json',
        '.xml': 'xml',
        '.py': 'code',
        '.js': 'code',
        '.ts': 'code',
        '.java': 'code',
        '.cpp': 'code',
        '.c': 'code',
        '.h': 'code',
        '.sql': 'code',
        '.yaml': 'text',
        '.yml': 'text',
        '.ini': 'text',
        '.cfg': 'text',
        '.conf': 'text',
    }
    
    @classmethod
    def can_extract(cls, file_path: str) -> bool:
        """Check if file type is supported"""
        ext = Path(file_path).suffix.lower()
        return ext in cls.SUPPORTED_EXTENSIONS
    
    @classmethod
    def get_file_type(cls, file_path: str) -> Optional[str]:
        """Get the file type category"""
        ext = Path(file_path).suffix.lower()
        return cls.SUPPORTED_EXTENSIONS.get(ext)
    
    @classmethod
    async def extract(cls, file_path: str) -> Optional[ExtractedContent]:
        """
        Extract text from a file.
        
        Args:
            file_path: Path to the file
            
        Returns:
            ExtractedContent object or None if extraction fails
        """
        if not os.path.exists(file_path):
            logger.error(f"File not found: {file_path}")
            return None
        
        file_type = cls.get_file_type(file_path)
        if not file_type:
            logger.warning(f"Unsupported file type: {file_path}")
            return None
        
        try:
            if file_type == 'pdf':
                return await cls._extract_pdf(file_path)
            elif file_type == 'docx':
                return await cls._extract_docx(file_path)
            elif file_type == 'doc':
                return await cls._extract_doc(file_path)
            elif file_type in ('pptx', 'ppt'):
                return await cls._extract_pptx(file_path)
            elif file_type in ('text', 'markdown', 'csv', 'json', 'xml', 'code'):
                return await cls._extract_text(file_path)
            elif file_type == 'html':
                return await cls._extract_html(file_path)
            elif file_type == 'rtf':
                return await cls._extract_rtf(file_path)
            else:
                return await cls._extract_text(file_path)
        except Exception as e:
            logger.error(f"Failed to extract text from {file_path}: {e}")
            return None
    
    @classmethod
    async def _extract_pdf(cls, file_path: str) -> Optional[ExtractedContent]:
        """Extract text from PDF file"""
        try:
            import pypdf
            
            text_parts = []
            page_count = 0
            metadata = {}
            
            with open(file_path, 'rb') as f:
                reader = pypdf.PdfReader(f)
                page_count = len(reader.pages)
                
                # Extract metadata
                if reader.metadata:
                    metadata = {
                        'title': str(reader.metadata.get('/Title', '')),
                        'author': str(reader.metadata.get('/Author', '')),
                        'subject': str(reader.metadata.get('/Subject', '')),
                        'creator': str(reader.metadata.get('/Creator', '')),
                    }
                
                # Extract text from each page
                for page_num, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
            
            full_text = "\n\n".join(text_parts)
            full_text = cls._clean_text(full_text)
            
            return ExtractedContent(
                text=full_text,
                metadata=metadata,
                page_count=page_count,
                extraction_method='pypdf'
            )
            
        except ImportError:
            logger.warning("pypdf not installed, trying pdfminer")
            return await cls._extract_pdf_fallback(file_path)
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            return await cls._extract_pdf_fallback(file_path)
    
    @classmethod
    async def _extract_pdf_fallback(cls, file_path: str) -> Optional[ExtractedContent]:
        """Fallback PDF extraction using pdfminer"""
        try:
            from pdfminer.high_level import extract_text
            
            text = extract_text(file_path)
            text = cls._clean_text(text)
            
            return ExtractedContent(
                text=text,
                metadata={'extractor': 'pdfminer'},
                extraction_method='pdfminer'
            )
        except ImportError:
            logger.error("Neither pypdf nor pdfminer installed")
            return None
        except Exception as e:
            logger.error(f"PDF fallback extraction failed: {e}")
            return None
    
    @classmethod
    async def _extract_docx(cls, file_path: str) -> Optional[ExtractedContent]:
        """Extract text from DOCX file"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            
            # Extract metadata
            metadata = {}
            if doc.core_properties:
                metadata = {
                    'title': doc.core_properties.title or '',
                    'author': doc.core_properties.author or '',
                    'subject': doc.core_properties.subject or '',
                    'created': str(doc.core_properties.created) if doc.core_properties.created else '',
                }
            
            # Extract text from paragraphs
            text_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        text_parts.append(row_text)
            
            full_text = "\n\n".join(text_parts)
            full_text = cls._clean_text(full_text)
            
            return ExtractedContent(
                text=full_text,
                metadata=metadata,
                extraction_method='python-docx'
            )
            
        except ImportError:
            logger.error("python-docx not installed. Run: pip install python-docx")
            return None
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            return None
    
    @classmethod
    async def _extract_doc(cls, file_path: str) -> Optional[ExtractedContent]:
        """Extract text from DOC file (legacy format)"""
        try:
            # Try antiword first (Linux)
            import subprocess
            result = subprocess.run(
                ['antiword', file_path],
                capture_output=True,
                text=True,
                timeout=30
            )
            if result.returncode == 0:
                text = cls._clean_text(result.stdout)
                return ExtractedContent(
                    text=text,
                    metadata={'extractor': 'antiword'},
                    extraction_method='antiword'
                )
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass
        except Exception:
            pass
        
        # Fallback: try to read as binary and extract text
        try:
            with open(file_path, 'rb') as f:
                content = f.read()
                # Simple text extraction (not reliable for complex DOC files)
                text = content.decode('utf-8', errors='ignore')
                text = re.sub(r'[^\x20-\x7E\n\t]', '', text)
                text = cls._clean_text(text)
                
                return ExtractedContent(
                    text=text,
                    metadata={'extractor': 'binary_fallback', 'warning': 'May be incomplete'},
                    extraction_method='binary_fallback'
                )
        except Exception as e:
            logger.error(f"DOC extraction failed: {e}")
            return None
    
    @classmethod
    async def _extract_pptx(cls, file_path: str) -> Optional[ExtractedContent]:
        """Extract text from PowerPoint PPTX files"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)

            text_parts = []
            slide_count = 0

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_count += 1
                slide_text_parts = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text_parts.append(shape.text)

                    # Extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = ' | '.join(
                                cell.text.strip() for cell in row.cells if cell.text.strip()
                            )
                            if row_text:
                                slide_text_parts.append(row_text)

                if slide_text_parts:
                    slide_header = f"--- Slide {slide_num} ---"
                    text_parts.append(slide_header)
                    text_parts.extend(slide_text_parts)

            full_text = "\n\n".join(text_parts)
            full_text = cls._clean_text(full_text)

            return ExtractedContent(
                text=full_text,
                metadata={'slide_count': slide_count},
                page_count=slide_count,
                extraction_method='python-pptx'
            )

        except ImportError:
            logger.error("python-pptx not installed. Run: pip install python-pptx")
            return None
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            return None

    @classmethod
    async def _extract_text(cls, file_path: str) -> Optional[ExtractedContent]:
        """Extract text from plain text files"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text = f.read()
                        text = cls._clean_text(text)
                        
                        return ExtractedContent(
                            text=text,
                            metadata={'encoding': encoding},
                            extraction_method='plaintext'
                        )
                except UnicodeDecodeError:
                    continue
            
            logger.error(f"Could not decode text file: {file_path}")
            return None
            
        except Exception as e:
            logger.error(f"Text extraction failed: {e}")
            return None
    
    @classmethod
    async def _extract_html(cls, file_path: str) -> Optional[ExtractedContent]:
        """Extract text from HTML file"""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                
                # Remove script and style elements
                for script in soup(["script", "style", "meta", "link"]):
                    script.decompose()
                
                # Get text
                text = soup.get_text(separator='\n')
                text = cls._clean_text(text)
                
                # Extract metadata
                metadata = {}
                title_tag = soup.find('title')
                if title_tag:
                    metadata['title'] = title_tag.get_text()
                
                return ExtractedContent(
                    text=text,
                    metadata=metadata,
                    extraction_method='beautifulsoup'
                )
                
        except ImportError:
            # Fallback: simple regex-based extraction
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                # Remove tags
                text = re.sub(r'<[^>]+>', ' ', content)
                text = cls._clean_text(text)
                
                return ExtractedContent(
                    text=text,
                    metadata={'extractor': 'regex_fallback'},
                    extraction_method='regex'
                )
        except Exception as e:
            logger.error(f"HTML extraction failed: {e}")
            return None
    
    @classmethod
    async def _extract_rtf(cls, file_path: str) -> Optional[ExtractedContent]:
        """Extract text from RTF file"""
        try:
            from striprtf.striprtf import rtf_to_text
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
                text = rtf_to_text(rtf_content)
                text = cls._clean_text(text)
                
                return ExtractedContent(
                    text=text,
                    metadata={'extractor': 'striprtf'},
                    extraction_method='striprtf'
                )
        except ImportError:
            # Fallback: simple RTF parsing
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                # Remove RTF control words
                text = re.sub(r'\\[a-z]+\d*\s?', '', content)
                text = re.sub(r'[{}]', '', text)
                text = cls._clean_text(text)
                
                return ExtractedContent(
                    text=text,
                    metadata={'extractor': 'regex_fallback'},
                    extraction_method='regex'
                )
        except Exception as e:
            logger.error(f"RTF extraction failed: {e}")
            return None
    
    @classmethod
    def _clean_text(cls, text: str) -> str:
        """Clean and normalize extracted text"""
        if not text:
            return ""
        
        # Normalize line endings
        text = re.sub(r'\r\n', '\n', text)
        text = re.sub(r'\r', '\n', text)
        
        # Remove excessive newlines (more than 2 consecutive)
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        # Remove excessive spaces
        text = re.sub(r'[ \t]+', ' ', text)
        
        # Remove null bytes and other control characters
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
        
        # Trim whitespace from each line
        lines = [line.strip() for line in text.split('\n')]
        text = '\n'.join(lines)
        
        # Final trim
        text = text.strip()
        
        return text


class TextChunker:
    """
    Splits text into semantic chunks for indexing.
    
    Uses a combination of:
    - Sentence boundaries
    - Paragraph boundaries  
    - Token limits
    - Overlap for context continuity
    """
    
    DEFAULT_CHUNK_SIZE = 500  # tokens (approximately)
    DEFAULT_CHUNK_OVERLAP = 50  # tokens
    CHARS_PER_TOKEN = 4  # Approximate for English text
    MIN_CHUNK_SIZE = 100  # Don't create tiny chunks
    
    @classmethod
    def chunk_text(
        cls,
        text: str,
        chunk_size: int = DEFAULT_CHUNK_SIZE,
        chunk_overlap: int = DEFAULT_CHUNK_OVERLAP,
        respect_paragraphs: bool = True
    ) -> List[Dict[str, Any]]:
        """
        Split text into overlapping chunks.
        
        Args:
            text: The text to chunk
            chunk_size: Target chunk size in tokens
            chunk_overlap: Overlap between chunks in tokens
            respect_paragraphs: Try to keep paragraphs together
            
        Returns:
            List of chunk dictionaries with text and metadata
        """
        if not text or not text.strip():
            return []
        
        # Convert token sizes to character sizes
        char_chunk_size = chunk_size * cls.CHARS_PER_TOKEN
        char_overlap = chunk_overlap * cls.CHARS_PER_TOKEN
        min_chars = cls.MIN_CHUNK_SIZE * cls.CHARS_PER_TOKEN
        
        chunks = []
        
        if respect_paragraphs:
            # First split by paragraphs
            paragraphs = cls._split_paragraphs(text)
            chunks = cls._chunk_by_paragraphs(
                paragraphs, char_chunk_size, char_overlap, min_chars
            )
        else:
            # Direct sentence-based chunking
            sentences = cls._split_sentences(text)
            chunks = cls._chunk_by_sentences(
                sentences, char_chunk_size, char_overlap, min_chars
            )
        
        # Add token counts and indices
        for i, chunk in enumerate(chunks):
            chunk['index'] = i
            chunk['token_count'] = len(chunk['text']) // cls.CHARS_PER_TOKEN
        
        return chunks
    
    @classmethod
    def _split_paragraphs(cls, text: str) -> List[str]:
        """Split text into paragraphs"""
        # Split on double newlines or more
        paragraphs = re.split(r'\n\s*\n', text)
        return [p.strip() for p in paragraphs if p.strip()]
    
    @classmethod
    def _split_sentences(cls, text: str) -> List[str]:
        """Split text into sentences"""
        # Simple sentence splitting on . ! ? followed by space/newline
        sentence_pattern = r'(?<=[.!?])\s+(?=[A-Z])'
        sentences = re.split(sentence_pattern, text)
        return [s.strip() for s in sentences if s.strip()]
    
    @classmethod
    def _chunk_by_paragraphs(
        cls,
        paragraphs: List[str],
        max_chars: int,
        overlap_chars: int,
        min_chars: int
    ) -> List[Dict[str, Any]]:
        """Chunk by keeping paragraphs together when possible"""
        chunks = []
        current_chunk_parts = []
        current_length = 0
        
        for para in paragraphs:
            para_length = len(para)
            
            # If single paragraph is too long, split it
            if para_length > max_chars:
                # Save current chunk first
                if current_chunk_parts:
                    chunk_text = '\n\n'.join(current_chunk_parts)
                    if len(chunk_text) >= min_chars:
                        chunks.append({'text': chunk_text})
                    current_chunk_parts = []
                    current_length = 0
                
                # Split long paragraph by sentences
                sentences = cls._split_sentences(para)
                sub_chunks = cls._chunk_by_sentences(
                    sentences, max_chars, overlap_chars, min_chars
                )
                chunks.extend(sub_chunks)
                continue
            
            # Check if adding this paragraph exceeds limit
            if current_length + para_length + 2 > max_chars and current_chunk_parts:
                # Save current chunk
                chunk_text = '\n\n'.join(current_chunk_parts)
                if len(chunk_text) >= min_chars:
                    chunks.append({'text': chunk_text})
                
                # Start new chunk - optionally with overlap
                if overlap_chars > 0 and current_chunk_parts:
                    # Use last paragraph as overlap if it's not too long
                    last_para = current_chunk_parts[-1]
                    if len(last_para) <= overlap_chars:
                        current_chunk_parts = [last_para]
                        current_length = len(last_para)
                    else:
                        current_chunk_parts = []
                        current_length = 0
                else:
                    current_chunk_parts = []
                    current_length = 0
            
            # Add paragraph to current chunk
            current_chunk_parts.append(para)
            current_length += para_length + 2  # +2 for \n\n
        
        # Don't forget the last chunk
        if current_chunk_parts:
            chunk_text = '\n\n'.join(current_chunk_parts)
            if len(chunk_text) >= min_chars:
                chunks.append({'text': chunk_text})
        
        return chunks
    
    @classmethod
    def _chunk_by_sentences(
        cls,
        sentences: List[str],
        max_chars: int,
        overlap_chars: int,
        min_chars: int
    ) -> List[Dict[str, Any]]:
        """Chunk by sentences with overlap"""
        chunks = []
        current_chunk = []
        current_length = 0
        
        for sentence in sentences:
            sentence_length = len(sentence)
            
            # If single sentence is too long, split it by words
            if sentence_length > max_chars:
                # Save current chunk first
                if current_chunk:
                    chunk_text = ' '.join(current_chunk)
                    if len(chunk_text) >= min_chars:
                        chunks.append({'text': chunk_text})
                    current_chunk = []
                    current_length = 0
                
                # Split long sentence
                for part in cls._split_long_text(sentence, max_chars):
                    if len(part) >= min_chars:
                        chunks.append({'text': part})
                continue
            
            # Check if adding sentence exceeds limit
            if current_length + sentence_length + 1 > max_chars and current_chunk:
                # Save current chunk
                chunk_text = ' '.join(current_chunk)
                if len(chunk_text) >= min_chars:
                    chunks.append({'text': chunk_text})
                
                # Start new chunk with overlap
                overlap_text = chunk_text[-overlap_chars:] if len(chunk_text) > overlap_chars else ""
                if overlap_text:
                    current_chunk = [overlap_text]
                    current_length = len(overlap_text)
                else:
                    current_chunk = []
                    current_length = 0
            
            # Add sentence to current chunk
            current_chunk.append(sentence)
            current_length += sentence_length + 1
        
        # Last chunk
        if current_chunk:
            chunk_text = ' '.join(current_chunk)
            if len(chunk_text) >= min_chars:
                chunks.append({'text': chunk_text})
        
        return chunks
    
    @classmethod
    def _split_long_text(cls, text: str, max_length: int) -> List[str]:
        """Split a long text into smaller parts at word boundaries"""
        parts = []
        words = text.split()
        current_part = []
        current_length = 0
        
        for word in words:
            word_len = len(word) + 1  # +1 for space
            
            if current_length + word_len > max_length and current_part:
                parts.append(' '.join(current_part))
                current_part = []
                current_length = 0
            
            current_part.append(word)
            current_length += word_len
        
        if current_part:
            parts.append(' '.join(current_part))
        
        return parts
